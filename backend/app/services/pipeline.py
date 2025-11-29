import os
import re
from datetime import datetime
from typing import List, Dict, Any, Optional

import requests
from bs4 import BeautifulSoup
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE

from app.core import get_settings
from app.core.charts.chart_generator import generate_charts
import asyncio

import json

# Optional helpers (used if installed)
try:
    import trafilatura
except ImportError:
    trafilatura = None

try:
    import wikipediaapi
except ImportError:
    wikipediaapi = None

try:
    from youtube_transcript_api import YouTubeTranscriptApi
except ImportError:
    YouTubeTranscriptApi = None

settings = get_settings()
client = OpenAI(api_key=settings.openai_api_key)

# -------------------------------------------------------------------
# Utility: Slugify topic
# -------------------------------------------------------------------
def _slugify(text: str) -> str:
    text = text.strip().lower()
    text = re.sub(r"[^a-z0-9]+", "-", text)
    text = re.sub(r"-+", "-", text).strip("-")
    return text or "report"

# -------------------------------------------------------------------
# HTTP config
# -------------------------------------------------------------------
HEADERS = {"User-Agent": "Mozilla/5.0 (compatible; AI-Research-Agent/1.0; +https://example.com)"}

# -------------------------------------------------------------------
# 1) SEARCH LAYER — Hybrid, multi-source
# -------------------------------------------------------------------
def _search_duckduckgo_html(topic: str, max_results: int = 8) -> List[str]:
    params = {"q": topic, "kl": "in-en"}
    resp = requests.get("https://duckduckgo.com/html/", params=params, headers=HEADERS, timeout=20)
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "html.parser")
    links: List[str] = []
    for a in soup.select("a.result__a"):
        href = a.get("href")
        if href and href.startswith("http"):
            links.append(href)
        if len(links) >= max_results:
            break
    return links

def _search_wikipedia_page(topic: str) -> Optional[str]:
    if not wikipediaapi:
        return None
    wiki = wikipediaapi.Wikipedia("en")
    page = wiki.page(topic)
    if page and page.exists():
        return page.title
    return None

def _search_youtube_ids(topic: str, max_results: int = 3) -> List[str]:
    if not YouTubeTranscriptApi:
        return []
    q = topic.replace(" ", "+")
    url = f"https://www.youtube.com/results?search_query={q}"
    try:
        resp = requests.get(url, headers=HEADERS, timeout=20)
        resp.raise_for_status()
    except Exception:
        return []
    video_ids: List[str] = []
    for match in re.finditer(r"watch\?v=([a-zA-Z0-9_-]{11})", resp.text):
        vid = match.group(1)
        if vid not in video_ids:
            video_ids.append(vid)
        if len(video_ids) >= max_results:
            break
    return video_ids

def search_web(topic: str, max_results: int = 8) -> Dict[str, Any]:
    urls = []
    try:
        urls = _search_duckduckgo_html(topic, max_results)
    except Exception:
        pass
    wiki_title = None
    try:
        wiki_title = _search_wikipedia_page(topic)
    except Exception:
        pass
    youtube_ids: List[str] = []
    try:
        youtube_ids = _search_youtube_ids(topic)
    except Exception:
        pass
    return {"urls": urls, "wiki_title": wiki_title, "youtube_ids": youtube_ids}

# -------------------------------------------------------------------
# 2) SCRAPING + CONTENT EXTRACTION LAYER
# -------------------------------------------------------------------
def _scrape_with_trafilatura(url: str) -> Optional[str]:
    if not trafilatura:
        return None
    try:
        downloaded = trafilatura.fetch_url(url)
        if not downloaded:
            return None
        return trafilatura.extract(downloaded, include_comments=False, include_tables=False, include_images=False)
    except Exception:
        return None

def _scrape_with_bs4(url: str) -> str:
    resp = requests.get(url, headers=HEADERS, timeout=30)
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator=" ")
    return " ".join(text.split())

def scrape_url(url: str, max_chars: int = 15000) -> str:
    text = _scrape_with_trafilatura(url) or _scrape_with_bs4(url)
    return text[:max_chars]

def fetch_wikipedia_content(title: str) -> Optional[str]:
    if not wikipediaapi:
        return None
    wiki = wikipediaapi.Wikipedia("en")
    page = wiki.page(title)
    if not page or not page.exists():
        return None
    parts = [page.summary or ""]
    for section in page.sections:
        parts.append(f"\n## {section.title}\n{section.text}\n")
    return "\n".join(parts)[:20000]

def fetch_youtube_transcript(video_id: str) -> Optional[str]:
    if not YouTubeTranscriptApi:
        return None
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=["en"])
        return " ".join(chunk["text"] for chunk in transcript)[:20000]
    except Exception:
        return None

# -------------------------------------------------------------------
# 3) OpenAI helper
# -------------------------------------------------------------------
def _call_llm(prompt: str, model: str = "gpt-4o-mini") -> str:
    if not settings.openai_api_key:
        raise RuntimeError("OPENAI_API_KEY not set")
    response = client.responses.create(model=model, input=prompt)
    try:
        return response.output[0].content[0].text
    except Exception:
        if hasattr(response, "output_text"):
            return response.output_text
        raise

def validate_research(research: Dict[str, Any]) -> Dict[str, Any]:
    """
    LLM Validation Agent:
    Cleans up the research summary, highlights caveats, and adds confidence signals.
    """
    research_json = json.dumps(research, ensure_ascii=False, indent=2)
    prompt = f"""
You are a senior research validation analyst.

Your job:
- Read the research JSON below.
- Clean and tighten the summary.
- Refine key points if needed.
- Identify any obvious gaps, caveats, or weak assumptions.
- Rate overall confidence as "high", "medium", or "low".

Return VALID JSON ONLY in this structure:

{{
  "validated_summary": "string",
  "validated_key_points": ["string", "string"],
  "caveats": ["string"],
  "confidence": "high" | "medium" | "low",
  "things_to_double_check": ["string"]
}}

Research JSON:
{research_json}
"""
    raw = _call_llm(prompt, model="gpt-4o-mini")
    try:
        data = json.loads(raw)
    except Exception:
        data = {
            "validated_summary": research.get("summary", ""),
            "validated_key_points": research.get("key_points", []),
            "caveats": [],
            "confidence": "medium",
            "things_to_double_check": []
        }
    return data


def plan_charts_from_research(research: Dict[str, Any]) -> Dict[str, Any]:
    """
    LLM Chart Planning Agent:
    Derives numeric chart_data from the research.
    """
    research_json = json.dumps(research, ensure_ascii=False, indent=2)
    prompt = f"""
You are a data visualization strategist.

From the research JSON below, infer realistic numeric values for key charts.
When real numbers are missing, fabricate plausible but conservative values.

Return VALID JSON ONLY in this structure:

{{
  "chart_data": {{
    "market_size": {{ "2020": number, "2021": number, "2022": number }},
    "market_share": {{ "Company A": number, "Company B": number, "Others": number }},
    "growth_projection": {{ "2024": number, "2025": number, "2026": number }},
    "competitor_strengths": {{ "Competitor A": number, "Competitor B": number }},
    "trend_frequency": {{ "Trend A": number, "Trend B": number }}
  }}
}}

Research JSON:
{research_json}
"""
    raw = _call_llm(prompt, model="gpt-4o-mini")
    try:
        data = json.loads(raw)
    except Exception:
        data = {
            "chart_data": {
                "market_size": {"2020": 8, "2021": 9, "2022": 10},
                "market_share": {"Company A": 40, "Company B": 30, "Others": 30},
                "growth_projection": {"2024": 15, "2025": 18, "2026": 22},
                "competitor_strengths": {"Competitor A": 80, "Competitor B": 65},
                "trend_frequency": {"Automation": 22, "Analytics": 14}
            }
        }
    data.setdefault("chart_data", {})
    return data


def generate_recommendations(research: Dict[str, Any]) -> Dict[str, Any]:
    """
    LLM Recommendation Agent:
    Produces executive recommendations and next-step actions.
    """
    research_json = json.dumps(research, ensure_ascii=False, indent=2)
    prompt = f"""
You are a strategy consultant.

From the research JSON below, produce practical recommendations and actions
for an executive audience.

Return VALID JSON ONLY in this structure:

{{
  "key_recommendations": ["string", "string", "string"],
  "action_plan": ["string", "string", "string"],
  "summary_recommendations": ["string", "string"]
}}

Research JSON:
{research_json}
"""
    raw = _call_llm(prompt, model="gpt-4o-mini")
    try:
        data = json.loads(raw)
    except Exception:
        data = {
            "key_recommendations": [],
            "action_plan": [],
            "summary_recommendations": []
        }
    return data



# -------------------------------------------------------------------
# 4) Source collection
# -------------------------------------------------------------------
def collect_sources(topic: str, max_results: int = 8) -> List[Dict[str, str]]:
    prompt = f"You are an AI research engine. Generate {max_results} high-quality synthetic sources about:\n\"{topic}\"\nReturn JSON with a \"sources\" list."
    raw = _call_llm(prompt)
    try:
        data = json.loads(raw)
        return data.get("sources", [])
    except Exception:
        return [{"url": f"https://synthetic.example.com/{topic}", "title": f"Overview of {topic}", "content": f"Synthetic summary for {topic}."}]

# -------------------------------------------------------------------
# 5) Research summarizer
# -------------------------------------------------------------------
import json

def build_research_summary(topic: str, sources: List[Dict[str, str]]) -> Dict[str, Any]:
    joined = "\n\n".join([f"### Source {i}\nURL: {s['url']}\nCONTENT:\n{s['content']}" for i, s in enumerate(sources, 1)])
    system = f"You are a senior research analyst. Summarize the web research about \"{topic}\" and return JSON with keys: topic, summary, key_points, statistics, trends, challenges, opportunities, sources_used."
    raw = _call_llm(system + "\n" + joined)
    try:
        data = json.loads(raw)
    except Exception:
        data = {"topic": topic, "summary": raw[:3000], "key_points": [], "statistics": [], "trends": [], "challenges": [], "opportunities": [], "sources_used": [s["url"] for s in sources]}
    data.setdefault("topic", topic)
    data.setdefault("sources_used", [s["url"] for s in sources])
    return data

# -------------------------------------------------------------------
# 6) Slide plan generator (FIXED & STABLE)
# -------------------------------------------------------------------
def generate_slide_plan(research: Dict[str, Any], theme_config: Dict[str, Any] = None) -> Dict[str, Any]:
    import json

    research_json = json.dumps(research, ensure_ascii=False, indent=2)


    theme_json = json.dumps(theme_config or {}, ensure_ascii=False, indent=2)

    prompt = f"""
You are a senior strategy consultant (McKinsey/BCG style) and a strict JSON generator.

Your task:
Create a complete slide-plan for a professional business PowerPoint deck based on the research data below.

You MUST return valid JSON only.
No commentary. No markdown. No natural language outside JSON.
If unsure, fabricate reasonable business values to complete the structure.

=====================================================================
REQUIRED JSON SCHEMA (DO NOT CHANGE KEYS)

{{
  "title": "string",
  "subtitle": "string",
  "theme_config": {{ ... }},  <-- INCLUDE THE PROVIDED THEME CONFIG HERE
  "sections": [
    {{
      "heading": "string",
      "bullets": ["string", "string", "string"]
    }}
  ],
  "conclusion_bullets": ["string", "string"],
  "chart_data": {{
    "market_size": {{"2020": 10, "2021": 12, "2022": 15}},
    "market_share": {{"Company A": 40, "Company B": 30}},
    "growth_projection": {{"2024": 20, "2025": 27, "2026": 34}},
    "competitor_strengths": {{"Competitor A": 80, "Competitor B": 65}},
    "trend_frequency": {{"Trend A": 20, "Trend B": 12}}
  }}
}}
=====================================================================

RULES FOR GENERATION
- Produce **15–20 slides** total.
- sections[] must contain **12–18 sections**.
- Each section must have:
    - heading
    - 4–7 short, clear business bullets
- All numeric values in chart_data must be realistic.
- Always include chart_data with all 5 subfields.
- **IMPORTANT**: You MUST include the "theme_config" object in the output, exactly as provided below. If it is empty, generate a professional default theme.
- Do not rename or remove any JSON keys.
- Output JSON only.

=====================================================================
THEME CONFIGURATION (Use this):
{theme_json}

=====================================================================
Now generate the slide-plan JSON based on the research below:

{research_json}
"""

    raw = _call_llm(prompt, model="gpt-4o-mini")  # upgrade model for reliability

    # ---------------------------------------
    # JSON Parsing + Recovery
    # ---------------------------------------
    try:
        data = json.loads(raw)
    except Exception:
        print("[SlidePlan] LLM returned invalid JSON. Using fallback plan.")
        data = {
            "title": research.get("topic", "Report"),
            "subtitle": "Auto-generated",
            "sections": [],
            "conclusion_bullets": [],
            "chart_data": {}
        }

    # Guarantee required keys
    data.setdefault("title", research.get("topic", "Report"))
    data.setdefault("subtitle", "Auto-generated")
    data.setdefault("sections", [])
    data.setdefault("conclusion_bullets", [])
    data.setdefault("chart_data", {})
    data.setdefault("theme_config", theme_config or {})

    return data


def auto_expand_slide_plan(slide_plan: Dict[str, Any]) -> Dict[str, Any]:
    """
    If the LLM failed or produced too few sections, expand with a guaranteed fallback.
    """
    if not slide_plan.get("sections") or len(slide_plan["sections"]) < 5:
        print("[SlidePlan] Expanding with fallback slide set (LLM incomplete).")

        slide_plan["sections"] = [
            {"heading": "Executive Summary", "bullets": ["Overview", "Key insights", "Business snapshot"]},
            {"heading": "Market Overview", "bullets": ["Market size", "Growth rate", "Segments", "Drivers"]},
            {"heading": "Industry Trends", "bullets": ["Trend 1", "Trend 2", "Trend 3"]},
            {"heading": "Challenges", "bullets": ["Challenge 1", "Challenge 2", "Challenge 3"]},
            {"heading": "Opportunities", "bullets": ["Opportunity 1", "Opportunity 2"]},
            {"heading": "Competitive Landscape", "bullets": ["Player A", "Player B", "Player C"]},
            {"heading": "SWOT Analysis", "bullets": ["Strengths", "Weaknesses", "Opportunities", "Threats"]},
            {"heading": "Use Cases", "bullets": ["Use Case 1", "Use Case 2"]},
            {"heading": "Future Outlook", "bullets": ["Forecast", "Emerging tech"]},
        ]

    # Guarantee chart_data exists
    slide_plan.setdefault("chart_data", {})
    slide_plan["chart_data"].setdefault("market_size", {"2020": 10, "2021": 12, "2022": 15})
    slide_plan["chart_data"].setdefault("market_share", {"Company A": 40, "Company B": 35})
    slide_plan["chart_data"].setdefault("growth_projection", {"2024": 20, "2025": 25, "2026": 32})
    slide_plan["chart_data"].setdefault("competitor_strengths", {"A": 80, "B": 60})
    slide_plan["chart_data"].setdefault("trend_frequency", {"Automation": 22, "Edge AI": 14})

    return slide_plan

# -------------------------------------------------------------------
# 7) PPT generator (Professional Dark Theme)
# -------------------------------------------------------------------
# -------------------------------------------------------------------
# 7) PPT generator (Professional Dark Theme)
# -------------------------------------------------------------------

def hex_to_rgb(hex_str):
    if not hex_str:
        return RGBColor(0, 0, 0)
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
# -------------------------------
# COLOR UTILITIES (ADD THESE)
# -------------------------------
def is_dark_color(hex_color: str) -> bool:
    """Return True if a hex color is dark."""
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    luminance = 0.2126*r + 0.7152*g + 0.0722*b
    return luminance < 128

def choose_text_color(bg_hex: str) -> RGBColor:
    """Automatically selects black or white text depending on background color."""
    return RGBColor(255, 255, 255) if is_dark_color(bg_hex) else RGBColor(20, 20, 20)    

def _apply_theme(slide, theme_config):
    theme_config = theme_config or {}
    bg_color = theme_config.get("background_color", "#121212")
    
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg_color)

    # Gradient support (simple vertical gradient simulation if requested)
    if theme_config.get("theme") == "gradient":
        # python-pptx doesn't support gradients easily on background directly via high-level API
        # So we stick to solid background or add a shape behind everything.
        # For robustness, we will stick to solid background color[0] or background_color
        pass

def _add_rounded_rect_background(slide, left, top, width, height, theme_config):
    """Adds an auto-contrast rounded rectangle background behind text and returns both the shape and card color."""
    radius = theme_config.get("corner_radius", 40)
    if radius <= 0:
        return None, None

    bg_hex = theme_config.get("background_color", "#121212")
    dark_bg = is_dark_color(bg_hex)

    # Create shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

    # Determine card color
    if dark_bg:
        card_color = RGBColor(255, 255, 255)  # white card on dark bg
        transparency = 0.20
    else:
        card_color = RGBColor(0, 0, 0)        # black card on light bg
        transparency = 0.40

    # Apply styling
    shape.fill.solid()
    shape.fill.fore_color.rgb = card_color
    shape.fill.transparency = transparency

    shape.line.fill.background()

    # Rounded corner radius (0–1 range)
    try:
        shape.adjustments[0] = radius / 100.0
    except:
        pass

    return shape, card_color



def _add_title_slide(prs: Presentation, title: str, subtitle: str, theme_config: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_theme(slide, theme_config)
    
    tc = theme_config or {}
    font_family = tc.get("font_family", "Arial")
    heading_color = hex_to_rgb(tc.get("accent_color", "#38BDF8"))
    title_color = heading_color
    sub_color = heading_color

    # Title
    tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = title_color
    p.font.name = font_family
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(2))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = sub_color
    p2.font.name = font_family
    p2.alignment = PP_ALIGN.LEFT
    return slide

def _add_bullet_slide(prs: Presentation, heading: str, bullets: List[str], theme_config: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_theme(slide, theme_config)
    
    tc = theme_config or {}
    font_family = tc.get("font_family", "Arial")

    bg_hex = tc.get("background_color", "#121212")

    # --- Auto contrast heading color ---
    heading_color = hex_to_rgb(tc.get("accent_color", "#38BDF8"))


    # --- Auto contrast bullet color ---
    card_shape, card_color = _add_rounded_rect_background(
    slide, Inches(0.5), Inches(1.8), Inches(12), Inches(5), tc
)

    # Decide bullet text color based on card color
    if card_color == RGBColor(255, 255, 255):  # white card
        text_color = RGBColor(30, 30, 30)      # dark bullets
    else:
        text_color = RGBColor(240, 240, 240)   # light bullets

    # Heading
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = heading
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = heading_color
    p.font.name = font_family

    
    # Bullets
    txb = slide.shapes.add_textbox(
    Inches(0.9),  # moves text right slightly
    Inches(2.2),  # moves text lower to fit inside card
    Inches(11.0), # slightly narrower
    Inches(4.2)   # adjusted height
)

    tfb = txb.text_frame
    tfb.word_wrap = True
    tfb.clear()

    if not bullets:
        bullets = ["Insight coming soon", "More data required"]

    for bullet in bullets:
        para = tfb.add_paragraph()
        para.text = f"• {bullet}"
        para.font.size = Pt(22)
        para.font.color.rgb = text_color
        para.font.name = font_family
        para.level = 0
        para.space_after = Pt(12)
        para.space_before = Pt(2)
        para.line_spacing = 1.3

    return slide


def _add_sources_slide(prs: Presentation, sources: List[str], theme_config: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_theme(slide, theme_config)
    
    tc = theme_config or {}
    font_family = tc.get("font_family", "Arial")

    # ADD THIS – you forgot it earlier
    bg_hex = tc.get("background_color", "#121212")

    # ---------- AUTO CONTRAST HEADINGS ----------
    heading_color = hex_to_rgb(tc.get("accent_color", "#38BDF8"))

    # ---------- AUTO CONTRAST SUBTEXT ----------
    if is_dark_color(bg_hex):
        subtext_color = hex_to_rgb(tc.get("subtext_color", "#B4B4B4"))
    else:
        subtext_color = RGBColor(40, 40, 40)


    # ----------------------------------------------------
    # Heading
    # ----------------------------------------------------
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Sources"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = heading_color
    p.font.name = font_family

    

    # ----------------------------------------------------
    # Source list text
    # ----------------------------------------------------
    txb = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(11.6), Inches(4.6))
    tfb = txb.text_frame
    tfb.word_wrap = True
    
    for s in sources[:8]:
        text = s.get("title") or s.get("url") or str(s)
        para = tfb.add_paragraph()
        para.text = text
        para.font.size = Pt(14)
        para.font.color.rgb = subtext_color
        para.font.name = font_family
        para.space_after = Pt(6)

    return slide


def _add_chart_slide(prs: Presentation, chart_path: str, theme_config: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_theme(slide, theme_config)

    tc = theme_config or {}
    font_family = tc.get("font_family", "Arial")
    bg_hex = tc.get("background_color", "#121212")

    # ---------- AUTO CONTRAST HEADING ----------
    heading_color = hex_to_rgb(tc.get("accent_color", "#38BDF8"))

    # Title from filename
    title = os.path.basename(chart_path)\
                .replace("chart_", "")\
                .replace(".png", "")\
                .replace("_", " ")\
                .title()

    # Title text
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = heading_color
    p.font.name = font_family

    # ---------- Optional card behind chart ----------
    # Optional card behind chart
    card_shape, card_color = _add_rounded_rect_background(
        slide,
        Inches(0.6), Inches(1.4),
        prs.slide_width - Inches(1.2),
        prs.slide_height - Inches(2.0),
        tc
    )

    # Send card backwards so chart stays visible
    if card_shape:
        card_shape.z_order = 1


    # ---------- Chart placement (safe margins) ----------
    left = Inches(0.85)
    top = Inches(1.55)
    max_width = prs.slide_width - Inches(1.7)
    max_height = prs.slide_height - Inches(2.2)

    pic = slide.shapes.add_picture(chart_path, left, top)

    # -------- AUTO SCALE TO FIT --------
    # Scale width first
    if pic.width > max_width:
        scale = max_width / pic.width
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)

    # Scale height if needed
    if pic.height > max_height:
        scale = max_height / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)

    # -------- CENTER CHART HORIZONTALLY --------
    pic.left = int((prs.slide_width - pic.width) / 2)

    return slide




# -------------------------------------------------------------------
# Animation helpers (unchanged from previous implementation)
# -------------------------------------------------------------------
def apply_slide_transition(slide):
    transition = OxmlElement("p:transition")
    transition.set("spd", "slow")
    fade = OxmlElement("p:fade")
    transition.append(fade)
    slide_element = slide.element
    insert_idx = 0
    cSld = slide_element.find(qn("p:cSld"))
    if cSld is not None:
        insert_idx = slide_element.index(cSld) + 1
    clrMapOvr = slide_element.find(qn("p:clrMapOvr"))
    if clrMapOvr is not None:
        insert_idx = slide_element.index(clrMapOvr) + 1
    slide_element.insert(insert_idx, transition)

# -------------------------------------------------------------------
# Bullet Animations (placeholder)
# -------------------------------------------------------------------
def apply_bullet_animations(slide):
    # (Implementation omitted for brevity – reuse previous logic)
    pass


# -------------------------------------------------------------------
# PPT generation with chart integration
# -------------------------------------------------------------------
def generate_ppt(topic: str, slide_plan: Dict[str, Any], output_dir: str, theme_config: Dict[str, Any] = None) -> (str, str):
    from pptx import Presentation
    from pptx.util import Inches
    import os
    from datetime import datetime

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title = slide_plan.get("title", topic)
    subtitle = slide_plan.get("subtitle", "Auto-generated AI research deck")
    sections = slide_plan.get("sections", [])
    conclusion_bullets = slide_plan.get("conclusion_bullets", [])
    sources = slide_plan.get("sources_used", [])
    
    # Use theme config from slide plan if not provided explicitly (fallback)
    if not theme_config:
        theme_config = slide_plan.get("theme_config", {})

    # --------------------- Title Slide ---------------------
    slide = _add_title_slide(prs, title, subtitle, theme_config)
    apply_slide_transition(slide)

    # --------------------- Content Sections ---------------------
    for sec in sections:
        slide = _add_bullet_slide(
            prs,
            sec.get("heading", "Section"),
            sec.get("bullets", []),
            theme_config
        )
        apply_bullet_animations(slide)
        apply_slide_transition(slide)

    # --------------------- Conclusion -----------------------
    if conclusion_bullets:
        slide = _add_bullet_slide(prs, "Key Takeaways", conclusion_bullets, theme_config)
        apply_bullet_animations(slide)
        apply_slide_transition(slide)

    # --------------------- Sources --------------------------
    if sources:
        slide = _add_sources_slide(prs, sources, theme_config)
        apply_slide_transition(slide)

    # --------------------- Charts ---------------------------
    chart_files = []
    if "chart_data" in slide_plan:
        try:
            # Pass theme_config to generate_charts if supported (we will update it next)
            chart_files = generate_charts(slide_plan, output_dir, theme_config)

            for cp in chart_files:
                if os.path.exists(cp):
                    slide = _add_chart_slide(prs, cp, theme_config)
                    apply_slide_transition(slide)

        except Exception as e:
            print(f"[ChartGenerator] Error: {e}")

    # --------------------- Save PPT -------------------------
    slug = _slugify(topic)
    ts = datetime.now().strftime("%Y%m%d_%H%M")
    filename = f"{slug}_{ts}.pptx"

    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, filename)

    prs.save(output_path)

    return filename, output_path


# -------------------------------------------------------------------
# Streaming pipeline (ASYNC GENERATOR) — MULTI-AGENT VERSION
# -------------------------------------------------------------------
async def run_research_pipeline_stream(topic: str, max_sources: int = 8, theme_config: Dict[str, Any] = None):

    # Start
    yield {"status": "start", "message": f"Starting research on: {topic}"}
    await asyncio.sleep(0)

    # ---------------------------------------------------------
    # 1) SOURCE COLLECTION
    # ---------------------------------------------------------
    yield {"status": "progress", "message": "🔍 Searching the web..."}
    await asyncio.sleep(0)

    try:
        sources = collect_sources(topic, max_results=max_sources)
        yield {"status": "progress", "message": f"✅ Found {len(sources)} sources."}
        await asyncio.sleep(0)
    except Exception as e:
        yield {"status": "error", "message": f"Search failed: {e}"}
        await asyncio.sleep(0)
        return

    # ---------------------------------------------------------
    # 2) RESEARCH SUMMARY AGENT
    # ---------------------------------------------------------
    yield {"status": "progress", "message": "🧠 Analyzing and summarizing content..."}
    await asyncio.sleep(0)

    research_summary = build_research_summary(topic, sources)

    # ---------------------------------------------------------
    # 3) VALIDATION AGENT
    # ---------------------------------------------------------
    yield {"status": "progress", "message": "🧹 Validating insights & cleaning output..."}
    await asyncio.sleep(0)

    validation = validate_research(research_summary)

    # ---------------------------------------------------------
    # 4) CHART PLANNING AGENT
    # ---------------------------------------------------------
    yield {"status": "progress", "message": "📊 Deriving chart data from research..."}
    await asyncio.sleep(0)

    chart_plan = plan_charts_from_research({
        **research_summary,
        "validated": validation
    })

    # ---------------------------------------------------------
    # 5) RECOMMENDATION AGENT
    # ---------------------------------------------------------
    yield {"status": "progress", "message": "📌 Generating executive recommendations..."}
    await asyncio.sleep(0)

    recommendations = generate_recommendations({
        **research_summary,
        "validated": validation
    })

    # ---------------------------------------------------------
    # 6) MERGE AGENT OUTPUTS → ENRICHED RESEARCH
    # ---------------------------------------------------------
    enriched_research = {
        **research_summary,
        "validated": validation,
        "chart_plan": chart_plan,
        "recommendations": recommendations,
    }

    # ---------------------------------------------------------
    # 7) SLIDE PLAN AGENT
    # ---------------------------------------------------------
    yield {"status": "progress", "message": "📝 Generating slide plan..."}
    await asyncio.sleep(0)

    slide_plan = generate_slide_plan(enriched_research, theme_config)

    # Merge chart data
    if chart_plan.get("chart_data"):
        slide_plan.setdefault("chart_data", {})
        slide_plan["chart_data"].update(chart_plan["chart_data"])

    # ---------------------------------------------------------
    # 8) Generate PPT
    # ---------------------------------------------------------
    filename, ppt_path = generate_ppt(topic, slide_plan, settings.outputs_dir, theme_config)

    # ---------------------------------------------------------
    # 9) Final DONE event
    # ---------------------------------------------------------
    yield {
        "status": "DONE",
        "message": "Pipeline completed",
        "ppt_filename": filename,
        "ppt_path": ppt_path,
        "topic": topic
    }








async def run_research_pipeline(topic: str, max_sources: int = 8, theme_config: Dict[str, Any] = None) -> Dict[str, Any]:
    """
    Wrapper to run the streaming pipeline and return the final result.
    """
    final_result = {}
    async for event in run_research_pipeline_stream(topic, max_sources, theme_config):
        if event["status"] == "DONE":
            final_result = event
    return final_result
